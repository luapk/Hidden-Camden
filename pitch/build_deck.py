"""
Camden Crawl partner deck — generates camden-crawl-deck.pptx
Run: python3 pitch/build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── colours ──────────────────────────────────────────────────────────────────
ACID       = RGBColor(0xCC, 0xFF, 0x00)   # #CCFF00
NIGHT      = RGBColor(0x00, 0x00, 0x00)   # pure black
NIGHT1     = RGBColor(0x0A, 0x0A, 0x0A)
NIGHT2     = RGBColor(0x11, 0x11, 0x13)
NIGHT3     = RGBColor(0x16, 0x16, 0x1A)
LABEL1     = RGBColor(0xF5, 0xF5, 0xF7)
LABEL2     = RGBColor(0x98, 0x98, 0x9D)
LABEL3     = RGBColor(0x5A, 0x5A, 0x5F)
CAMDEN_RED = RGBColor(0xD8, 0x43, 0x2F)
LINE       = RGBColor(0x28, 0x28, 0x2C)   # approx white/10%

# ── slide size: 16:9 ─────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

# ── helper primitives ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text(slide, text, x, y, w, h,
             font_name="Jost", size=18, bold=False,
             color=LABEL1, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def kicker(slide, text, x, y, w=Inches(10), color=ACID):
    add_text(slide, text.upper(), x, y, w, Inches(0.35),
             font_name="Space Grotesk", size=9, bold=True,
             color=color, align=PP_ALIGN.LEFT)

def display_head(slide, text, x, y, w, size=48, color=LABEL1):
    add_text(slide, text.upper(), x, y, w, Inches(2.5),
             font_name="Anton", size=size, bold=False,
             color=color, align=PP_ALIGN.LEFT)

def body(slide, text, x, y, w, h=Inches(1.2), size=14, color=LABEL2):
    add_text(slide, text, x, y, w, h,
             font_name="Jost", size=size, color=color)

def stat_block(slide, num, label, x, y):
    add_text(slide, num, x, y, Inches(3.5), Inches(1),
             font_name="Anton", size=52, color=ACID)
    add_text(slide, label, x, y + Inches(1), Inches(3.5), Inches(0.8),
             font_name="Space Grotesk", size=11, color=LABEL2)

def divider(slide, x, y, w):
    add_rect(slide, x, y, w, Pt(1.5), LINE)

def footer(slide, page, total=10):
    # brandmark left
    add_text(slide, "CAMDEN CRAWL", Inches(0.5), H - Inches(0.45),
             Inches(3), Inches(0.3), font_name="Space Grotesk",
             size=8, color=LABEL3, bold=True)
    # page num right
    add_text(slide, f"{page:02d} / {total:02d}", W - Inches(1.5), H - Inches(0.45),
             Inches(1.4), Inches(0.3), font_name="Space Grotesk",
             size=8, color=LABEL3, align=PP_ALIGN.RIGHT)

def bg(slide, color=NIGHT):
    add_rect(slide, 0, 0, W, H, color)

def placeholder_img(slide, x, y, w, h, label):
    """Dark crosshatch box with an acid-green label."""
    add_rect(slide, x, y, w, h, NIGHT2)
    # diagonal lines via two rectangles (pptx can't do patterns easily)
    add_rect(slide, x, y, w, Pt(1), LINE)
    add_rect(slide, x, y+h-Pt(1), w, Pt(1), LINE)
    add_text(slide, f"[ IMAGE · {label.upper()} ]",
             x + Inches(0.2), y + h//2 - Inches(0.25),
             w - Inches(0.4), Inches(0.5),
             font_name="Space Grotesk", size=9,
             color=RGBColor(0x66, 0x88, 0x00), align=PP_ALIGN.CENTER)

def acid_pill(slide, text, x, y):
    """Acid-green filled label chip."""
    box = slide.shapes.add_shape(1, x, y, Inches(1.8), Inches(0.32))
    box.fill.solid(); box.fill.fore_color.rgb = ACID
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.name = "Space Grotesk"; r.font.size = Pt(8.5)
    r.font.bold = True; r.font.color.rgb = NIGHT

def dark_card(slide, x, y, w, h):
    r = add_rect(slide, x, y, w, h, NIGHT2)
    # thin border
    r.line.color.rgb = LINE
    r.line.width = Pt(0.75)
    return r

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
placeholder_img(s, 0, 0, W, H, "Camden High Street at night · neon · crowd")
# dark gradient overlay rectangle (left half)
add_rect(s, 0, 0, Inches(8), H, NIGHT)
# make it semi-transparent via XML
overlay = s.shapes[-1]
sp = overlay.element
sp_pr = sp.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')

kicker(s, "Geofenced audio tours · Camden Town", Inches(0.85), Inches(0.85))
display_head(s, "Camden\nCrawl", Inches(0.85), Inches(1.4), Inches(8), size=96)
body(s, "The streets tell the stories. The bars pour the drinks.\nA walking tour that turns Camden's footfall into verified visits\nfor the brands that built it.",
     Inches(0.85), Inches(5.6), Inches(6.8), h=Inches(1.5), size=16, color=LABEL1)
footer(s, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — THE OPPORTUNITY
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "The Opportunity", Inches(0.85), Inches(0.72))
display_head(s, "Camden is one of the most walked corners of London.\nAlmost none of it is measurable.",
             Inches(0.85), Inches(1.15), Inches(11.5), size=36)
body(s, "Millions move through these streets every year. Brands pay for windows, billboards and pop-ups, then guess at the return. We make the walk itself the medium, and the redemption the receipt.",
     Inches(0.85), Inches(3.05), Inches(11.5), h=Inches(0.9), size=14, color=LABEL2)

divider(s, Inches(0.85), Inches(4.15), Inches(11.6))

stat_block(s, "28M",    "visitors to Camden Market\nevery year", Inches(0.85), Inches(4.4))
stat_block(s, "250K",   "people through Camden Town\nevery week, skewed under-35s", Inches(5.1),  Inches(4.4))
stat_block(s, "91%",    "of consumers feel more positive\nabout a brand after a live experience", Inches(9.35), Inches(4.4))
footer(s, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — WHY NOW (full bleed left)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
placeholder_img(s, 0, 0, Inches(5.5), H, "headphones · phone · Camden mural")
# right half content
add_rect(s, Inches(5.5), 0, Inches(7.83), H, NIGHT)
kicker(s, "Why Now", Inches(6.0), Inches(0.85))
display_head(s, "The world put\nits headphones on.", Inches(6.0), Inches(1.3), Inches(7.0), size=40)

bullets = [
    ("Location audio is intimate, hands-free and unskippable in a way a poster never is.", LABEL1),
    ("A two-minute story earns minutes of attention — not a 1.5-second glance.", LABEL1),
    ("Every stop ends with an action, not an impression.", LABEL1),
]
by = Inches(3.6)
for txt, col in bullets:
    add_text(s, "→", Inches(6.0), by, Inches(0.4), Inches(0.5),
             font_name="Jost", size=14, color=ACID, bold=True)
    add_text(s, txt, Inches(6.5), by, Inches(6.6), Inches(0.55),
             font_name="Jost", size=14, color=col)
    by += Inches(0.72)

footer(s, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — CUSTOMER JOURNEY
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "The Customer Journey", Inches(0.85), Inches(0.72))
display_head(s, "Walk. Unlock. Listen. Bank. Redeem.", Inches(0.85), Inches(1.15), Inches(11.5), size=38)
body(s, "One route from Camden Town tube. Each venue unlocks only when you physically arrive.\nNo tickets to print. No hardware in the bar. No app store.",
     Inches(0.85), Inches(2.65), Inches(11.5), h=Inches(0.75), size=14, color=LABEL2)

steps = [
    ("1", "Arrive",    "You reach a venue on the route.\nThe stop wakes up."),
    ("2", "Unlock",    "A two-minute story about\nthat spot plays in your ears."),
    ("3", "Bank",      "A sponsored drink lands in\nyour wallet for 7 days."),
    ("4", "Redeem",    "At the bar a live ticket counts\ndown. Staff tear it to pour."),
    ("5", "Walk On",   "The next story plays you\ndown to the next stop."),
]
sx = Inches(0.5)
sw = Inches(2.44)
gap = Inches(0.18)
sy = Inches(3.55)
for i, (num, head, desc) in enumerate(steps):
    dark_card(s, sx, sy, sw, Inches(3.25))
    # acid circle
    circ = s.shapes.add_shape(9, sx + Inches(0.16), sy + Inches(0.2), Inches(0.6), Inches(0.6))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACID
    circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap=False
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=num
    r.font.name="Anton"; r.font.size=Pt(20); r.font.color.rgb=NIGHT; r.font.bold=True
    add_text(s, head.upper(), sx + Inches(0.16), sy + Inches(0.95), sw - Inches(0.3), Inches(0.5),
             font_name="Anton", size=20, color=LABEL1)
    add_text(s, desc, sx + Inches(0.16), sy + Inches(1.5), sw - Inches(0.3), Inches(1.2),
             font_name="Jost", size=12, color=LABEL2)
    if i < 4:
        add_text(s, "→", sx + sw + Inches(0.02), sy + Inches(0.2), Inches(0.18), Inches(0.4),
                 font_name="Jost", size=16, color=ACID, bold=True)
    sx += sw + gap + Inches(0.18)

footer(s, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — BRAND TOUCHPOINTS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "Brand Touchpoints", Inches(0.85), Inches(0.72))
display_head(s, "Four places a brand shows up.", Inches(0.85), Inches(1.15), Inches(11.5), size=44)
body(s, "Not a logo in a corner. A presence woven through the walk, the story and the drink in their hand.",
     Inches(0.85), Inches(2.7), Inches(11.5), h=Inches(0.55), size=14, color=LABEL2)

cards = [
    ("In the ear",  "The Story",  "Your brand named in the narration of the stop it sponsors. Heritage, not advertising."),
    ("In the hand", "The Drink",  "A pint, a pour, a product sample. Banked and redeemed in venue. The payoff is yours."),
    ("On screen",   "The Ticket", "Co-branded redemption ticket at the moment of highest attention in the whole journey."),
    ("In store",    "The Stop",   "Your shop becomes a destination on the route. Footfall walked to your door."),
]
cx = Inches(0.5)
cw = Inches(3.0)
for tick, head, desc in cards:
    dark_card(s, cx, Inches(3.55), cw, Inches(3.25))
    add_text(s, tick.upper(), cx + Inches(0.22), Inches(3.77), cw - Inches(0.44), Inches(0.35),
             font_name="Space Grotesk", size=9, color=ACID, bold=True)
    add_text(s, head.upper(), cx + Inches(0.22), Inches(4.2), cw - Inches(0.44), Inches(0.75),
             font_name="Anton", size=30, color=LABEL1)
    add_text(s, desc, cx + Inches(0.22), Inches(5.05), cw - Inches(0.44), Inches(1.4),
             font_name="Jost", size=12.5, color=LABEL2)
    cx += cw + Inches(0.27)

footer(s, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — THE VOICES
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "The Voices", Inches(0.85), Inches(0.72))
display_head(s, "Camden, narrated by the people who made it.", Inches(0.85), Inches(1.15), Inches(11.5), size=38)
body(s, "Signature tours voiced by Camden's own. A name on the route sells the ticket and earns the press.",
     Inches(0.85), Inches(2.65), Inches(11.5), h=Inches(0.55), size=14, color=LABEL2)

celebs = [
    ("Suggs",       "Madness · the Dublin Castle years"),
    ("Yungblud",    "The new Camden, loud and now"),
    ("Carl Barât",  "The Libertines · the Albion circuit"),
]
cx = Inches(0.5)
cw = Inches(4.04)
for name, role in celebs:
    placeholder_img(s, cx, Inches(3.4), cw, Inches(3.42), f"portrait · {name}")
    # gradient overlay at bottom of portrait
    add_rect(s, cx, Inches(5.5), cw, Inches(1.32), NIGHT)
    add_text(s, name.upper(), cx + Inches(0.18), Inches(5.55), cw, Inches(0.75),
             font_name="Anton", size=34, color=LABEL1)
    add_text(s, role.upper(), cx + Inches(0.18), Inches(6.25), cw, Inches(0.35),
             font_name="Space Grotesk", size=9, color=ACID, bold=True)
    cx += cw + Inches(0.19)

footer(s, 6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — PARTNERS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "Built for partners like you", Inches(0.85), Inches(0.72))
display_head(s, "Three names that\nalready are Camden.", Inches(0.85), Inches(1.15), Inches(6), size=44)

partner_bullets = [
    ("Camden Hells",      "The drink in the wallet. Every banked pint is a verified pour, billed only when a real customer redeems in a real venue."),
    ("Dr. Martens",       "A stop on the route and a reward at the till. Walk the crowd to the store; hand them a reason to cross the threshold."),
    ("Universal Music",   "The new market store as a flagship stop. The tour scores the visit, the shop closes it."),
]
by = Inches(3.3)
for name, desc in partner_bullets:
    add_text(s, "→ " + name, Inches(0.85), by, Inches(5.8), Inches(0.45),
             font_name="Jost", size=15, bold=True, color=ACID)
    add_text(s, desc, Inches(1.15), by + Inches(0.4), Inches(5.5), Inches(0.7),
             font_name="Jost", size=13, color=LABEL2)
    by += Inches(1.1)

# right column
add_rect(s, Inches(7.3), Inches(1.1), Inches(5.7), H - Inches(1.55), NIGHT2)
kicker(s, "The shape of a partnership", Inches(7.65), Inches(1.4), color=LABEL3)

mini = [
    ("Reach",  "Every walker on the route, every day it runs."),
    ("Proof",  "A dashboard of verified redemptions, not estimated impressions."),
    ("Risk",   "Pay on result. No stock held, no wastage, no guesswork."),
]
my = Inches(1.95)
for m_head, m_body in mini:
    dark_card(s, Inches(7.65), my, Inches(5.1), Inches(1.15))
    add_text(s, m_head.upper(), Inches(7.9), my + Inches(0.14), Inches(4.8), Inches(0.35),
             font_name="Space Grotesk", size=9, color=ACID, bold=True)
    add_text(s, m_body, Inches(7.9), my + Inches(0.5), Inches(4.7), Inches(0.55),
             font_name="Jost", size=12.5, color=LABEL1)
    my += Inches(1.3)

footer(s, 7)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — THE CASE FOR BRANDS (full bleed right)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
placeholder_img(s, Inches(6.2), 0, Inches(7.13), H, "pint poured · neon · paper ticket on bar")
add_rect(s, 0, 0, Inches(6.5), H, NIGHT)

kicker(s, "The case for brands", Inches(0.85), Inches(0.85))
display_head(s, "You only pay\nwhen it works.", Inches(0.85), Inches(1.3), Inches(5.8), size=50)
body(s, "Traditional outdoor spend buys a guess. We bill against a redemption that happened, in a venue, by a real person who walked there. Every pour is a data point: who, where, when, how often.",
     Inches(0.85), Inches(3.5), Inches(5.6), h=Inches(1.2), size=13.5, color=LABEL2)

divider(s, Inches(0.85), Inches(4.9), Inches(5.6))
stat_block(s, "1:1",   "spend tied directly to\na verified redemption", Inches(0.85), Inches(5.1))
stat_block(s, "100%",  "activations location\nverified", Inches(2.8),   Inches(5.1))
stat_block(s, "0",     "hardware or wastage\nfor any venue", Inches(4.8),   Inches(5.1))

footer(s, 8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — PRICING MODELS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "Three ways to partner", Inches(0.85), Inches(0.72))
display_head(s, "Pick the model that fits the goal.", Inches(0.85), Inches(1.15), Inches(11.5), size=40)

models = [
    {
        "tag": "Model A",
        "head": "Pay Per Pour",
        "price": "£ / unit",
        "price_sub": "per verified redemption",
        "bullets": [
            "Pure performance — pay only when a drink or reward is claimed",
            "Lowest commitment, fastest to launch",
            "Best for testing Camden before scaling",
        ],
        "ideal": "Ideal: Camden Hells",
        "best": False,
    },
    {
        "tag": "Model B",
        "head": "Stop Sponsorship",
        "price": "£/month",
        "price_sub": "+ per redemption",
        "bullets": [
            "Own a stop: story, ticket and reward",
            "Monthly fee for presence plus a per-pour top-up",
            "Co-branded ticket and named narration",
        ],
        "ideal": "Ideal: Dr. Martens · Universal Music",
        "best": True,
    },
    {
        "tag": "Model C",
        "head": "Headline Activation",
        "price": "£ fixed",
        "price_sub": "campaign / launch",
        "bullets": [
            "A whole tour or celebrity route, brand-led",
            "Launch moment, PR and full data package",
            "Best for a store opening or product drop",
        ],
        "ideal": "Ideal: a flagship moment",
        "best": False,
    },
]

mx = Inches(0.5)
mw = Inches(4.04)
for m in models:
    dark_card(s, mx, Inches(2.65), mw, Inches(4.4))
    ty = Inches(2.85)
    if m["best"]:
        acid_pill(s, "Most Popular", mx + Inches(0.22), ty)
        ty += Inches(0.42)
    add_text(s, m["tag"].upper(), mx + Inches(0.22), ty, mw - Inches(0.44), Inches(0.32),
             font_name="Space Grotesk", size=9, color=ACID, bold=True)
    ty += Inches(0.38)
    add_text(s, m["head"].upper(), mx + Inches(0.22), ty, mw - Inches(0.44), Inches(0.65),
             font_name="Anton", size=30, color=LABEL1)
    ty += Inches(0.72)
    add_text(s, m["price"], mx + Inches(0.22), ty, mw - Inches(0.44), Inches(0.65),
             font_name="Anton", size=34, color=ACID)
    add_text(s, m["price_sub"], mx + Inches(0.22), ty + Inches(0.6), mw - Inches(0.44), Inches(0.32),
             font_name="Space Grotesk", size=10, color=LABEL2)
    ty += Inches(1.05)
    for b in m["bullets"]:
        add_text(s, "+ " + b, mx + Inches(0.22), ty, mw - Inches(0.44), Inches(0.52),
                 font_name="Jost", size=11.5, color=LABEL1)
        ty += Inches(0.56)
    add_text(s, m["ideal"], mx + Inches(0.22), Inches(6.7), mw - Inches(0.44), Inches(0.32),
             font_name="Space Grotesk", size=9, color=LABEL3)
    mx += mw + Inches(0.27)

footer(s, 9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — CLOSE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
placeholder_img(s, 0, 0, W, H, "Camden skyline dusk · lock bridge · acid light")
add_rect(s, 0, 0, Inches(8.5), H, NIGHT)

kicker(s, "Let's walk it", Inches(0.85), Inches(0.85))
display_head(s, "Be the reason\nthey keep\nwalking.", Inches(0.85), Inches(1.5), Inches(7.8), size=78)

# "walking." last word in acid
add_text(s, "WALKING.", Inches(0.85), Inches(4.8), Inches(7.8), Inches(1.1),
         font_name="Anton", size=78, color=ACID)

divider(s, Inches(0.85), Inches(5.9), Inches(5.5))
body(s, "A pilot route is live. We can have your brand on it\nin weeks, not quarters.",
     Inches(0.85), Inches(6.05), Inches(6.5), h=Inches(0.7), size=15, color=LABEL1)
add_text(s, "paulknott@gmail.com", Inches(0.85), Inches(6.85), Inches(6.5), Inches(0.4),
         font_name="Space Grotesk", size=13, color=ACID, bold=True)

footer(s, 10)

# ─────────────────────────────────────────────────────────────────────────────
out = "pitch/camden-crawl-deck.pptx"
prs.save(out)
print(f"Saved: {out}")
